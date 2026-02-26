"""
Tokuro AI Report Generation Engine
DAAKYI's Proprietary AI-Powered Report Generation System

Features:
- Multi-format report generation (PDF, PPTX, DOCX, XLSX)
- Tokuro AI-driven content intelligence
- Dynamic template processing
- Professional branding integration
- Executive-ready presentations
"""

import os
import json
import uuid
import asyncio
from typing import Dict, List, Optional, Any, Union
from datetime import datetime, timedelta
from pathlib import Path
import logging

# External libraries for report generation
from reportlab.lib.pagesizes import letter, A4
from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, Image, Table, TableStyle
from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
from reportlab.lib import colors
from reportlab.lib.units import inch
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from openpyxl import Workbook
from openpyxl.styles import Font, Alignment, PatternFill
from docx import Document
from docx.shared import Inches as DocxInches
from jinja2 import Template, Environment, FileSystemLoader
import aiofiles
import base64
from io import BytesIO

from models import (
    ReportGeneration, ReportTemplate, ReportType, ReportFormat, ReportStatus,
    TokuroAIReportData, User, Assessment, Organization
)
from database import DatabaseOperations
from tokuro_ai_engine import TokuroAIService

logger = logging.getLogger(__name__)

class TokuroReportEngine:
    """
    Tokuro AI-powered report generation engine
    Integrates with DAAKYI's proprietary AI for intelligent compliance reporting
    """
    
    def __init__(self):
        self.output_dir = Path("generated_reports")
        self.template_dir = Path("report_templates")
        self.assets_dir = Path("report_assets")
        
        # Ensure directories exist
        self.output_dir.mkdir(exist_ok=True)
        self.template_dir.mkdir(exist_ok=True)
        self.assets_dir.mkdir(exist_ok=True)
        
        # Initialize Tokuro AI engine
        self.tokuro_ai = TokuroAIService()
        
        # Initialize Jinja2 template environment
        self.jinja_env = Environment(
            loader=FileSystemLoader(str(self.template_dir))
        )
        
        # DAAKYI branding configuration
        self.branding = {
            'primary_color': '#1f2937',  # DAAKYI primary
            'secondary_color': '#3b82f6',  # DAAKYI secondary
            'accent_color': '#f59e0b',  # DAAKYI accent
            'tokuro_color': '#6b7280',  # Tokuro grey
            'logo_path': str(self.assets_dir / 'daakyi_logo.png'),
            'tokuro_logo_path': str(self.assets_dir / 'tokuro_ai_logo.png'),
            'watermark_text': 'Powered by Tokuro AI',
            'fonts': {
                'primary': 'Inter',
                'secondary': 'Nunito Sans'
            }
        }

    async def generate_report(
        self, 
        report_config: ReportGeneration, 
        assessment_data: Optional[Dict[str, Any]] = None
    ) -> str:
        """
        Generate a report based on configuration
        Returns the file path of the generated report
        """
        try:
            # Update status to generating
            await self._update_report_status(
                report_config.id, 
                ReportStatus.GENERATING, 
                progress=10
            )
            
            # Gather data for report
            logger.info(f"Gathering data for report {report_config.id}")
            report_data = await self._gather_report_data(report_config, assessment_data)
            await self._update_report_progress(report_config.id, 30)
            
            # Generate Tokuro AI insights if enabled
            if report_config.include_tokuro_ai:
                logger.info("Generating Tokuro AI insights")
                tokuro_data = await self._generate_tokuro_insights(report_data)
                report_config.tokuro_ai_data = tokuro_data
                await self._update_report_progress(report_config.id, 50)
            
            # Load template
            template = await self._load_template(report_config.template_id)
            await self._update_report_progress(report_config.id, 60)
            
            # Generate report based on format
            output_path = await self._generate_by_format(
                report_config, 
                report_data, 
                template
            )
            await self._update_report_progress(report_config.id, 90)
            
            # Finalize report
            await self._finalize_report(report_config, output_path)
            await self._update_report_status(
                report_config.id, 
                ReportStatus.COMPLETED, 
                progress=100
            )
            
            logger.info(f"Report {report_config.id} generated successfully: {output_path}")
            return output_path
            
        except Exception as e:
            logger.error(f"Report generation failed for {report_config.id}: {str(e)}")
            await self._update_report_status(
                report_config.id, 
                ReportStatus.FAILED, 
                error_message=str(e)
            )
            raise

    async def _gather_report_data(
        self,
        config: ReportGeneration,
        assessment_data: Optional[Dict[str, Any]] = None
    ) -> Dict[str, Any]:
        """Gather all necessary data for report generation.

        Data sources (in priority order):
          1. mvp1_intake_forms  — primary findings / gap source (engagement_id)
          2. mvp1_evidence      — evidence items for review counts (engagement_id)
          3. control_assessments / assessments — legacy collections (assessment_id)
        """

        data = {
            'report_info': {
                'id': config.id,
                'name': config.name,
                'description': config.description,
                'generated_at': datetime.utcnow(),
                'report_type': config.report_type,
                'format': config.report_format
            },
            'branding': self.branding,
            'organization': {},
            'assessment': {},
            'evidence_files': [],
            'control_assessments': [],
            'framework_data': {},
            'metadata': {},
            'generated_findings': [],
            'intake_forms': [],
            'mvp1_evidence': [],
        }

        # ── Organization ────────────────────────────────────────────────────────
        try:
            if config.organization_id:
                org = await DatabaseOperations.find_one(
                    "organizations",
                    {"id": config.organization_id}
                )
                if org:
                    data['organization'] = org
                    if 'branding_config' in org:
                        data['branding'].update(org['branding_config'])
        except Exception as org_err:
            logger.debug(f"Could not fetch organization: {org_err}")

        # ── Legacy assessments collection (optional) ────────────────────────────
        try:
            if config.assessment_id:
                assessment = await DatabaseOperations.find_one(
                    "assessments",
                    {"id": config.assessment_id}
                )
                if assessment:
                    data['assessment'] = assessment

                legacy_controls = await DatabaseOperations.find_many(
                    "control_assessments",
                    {"assessment_id": config.assessment_id}
                )
                data['control_assessments'] = legacy_controls or []

                legacy_evidence = await DatabaseOperations.find_many(
                    "evidence_files",
                    {"assessment_id": config.assessment_id}
                )
                data['evidence_files'] = legacy_evidence or []
        except Exception as leg_err:
            logger.debug(f"Legacy assessment collections not found or error: {leg_err}")

        if assessment_data:
            data['assessment'].update(assessment_data)

        # ── MVP1: mvp1_intake_forms (primary findings source) ───────────────────
        intake_forms: list = []
        try:
            intake_forms = await DatabaseOperations.find_many(
                "mvp1_intake_forms",
                {"engagement_id": config.assessment_id}
            ) or []
            data['intake_forms'] = intake_forms
            logger.info(
                f"[report {config.id}] mvp1_intake_forms: {len(intake_forms)} records "
                f"for engagement_id={config.assessment_id}"
            )
        except Exception as if_err:
            logger.warning(f"Could not query mvp1_intake_forms: {if_err}")

        # ── MVP1: mvp1_evidence (evidence items for review) ─────────────────────
        evidence_list: list = []
        try:
            evidence_list = await DatabaseOperations.find_many(
                "mvp1_evidence",
                {"engagement_id": config.assessment_id}
            ) or []
            data['mvp1_evidence'] = evidence_list
            logger.info(
                f"[report {config.id}] mvp1_evidence: {len(evidence_list)} records "
                f"for engagement_id={config.assessment_id}"
            )
        except Exception as ev_err:
            logger.warning(f"Could not query mvp1_evidence: {ev_err}")

        # ── Build analyst name lookup ────────────────────────────────────────────
        analyst_map: dict = {}
        try:
            analyst_ids = list({
                f.get('analyst_id') for f in intake_forms if f.get('analyst_id')
            })
            if analyst_ids:
                analysts = await DatabaseOperations.find_many(
                    "mvp1_users", {"id": {"$in": analyst_ids}}
                )
                analyst_map = {a['id']: a.get('name', 'Unknown Analyst') for a in (analysts or [])}
        except Exception as an_err:
            logger.debug(f"Analyst lookup failed: {an_err}")

        # ── Build generated_findings ─────────────────────────────────────────────
        #
        # Each finding carries three fields the PDF tables read directly:
        #   • status              — raw DB value (e.g. "gap_found", "ai_verified")
        #   • remediation_summary — the AI's reasoning / gap narrative
        #   • ai_confidence_score — 0.0–1.0 float (also stored as aiConfidence int %)
        #
        findings_results: list = []

        # 1. Intake forms → findings (primary source)
        for form in intake_forms:
            if not isinstance(form, dict):
                continue

            raw_status    = form.get('status') or 'draft'
            confidence    = float(form.get('ai_confidence_score') or 0.0)
            ai_results    = form.get('ai_analysis_results') or {}
            gaps          = form.get('identified_gaps') or []

            # Status resolution — honour DB value when it is a known terminal state;
            # otherwise derive from confidence / gaps so pending items still appear.
            TERMINAL_STATUSES = {
                'approved', 'rejected', 'ai_verified',
                'gap_found', 'needs_remediation', 'requires_revision', 'flagged',
            }
            if raw_status in TERMINAL_STATUSES:
                status_val = raw_status
            elif confidence >= 0.8 and not gaps:
                status_val = 'ai_verified'
            elif confidence > 0:
                status_val = 'pending'
                if not gaps:
                    gaps = ['Manual review recommended']
            else:
                status_val = 'pending_ai_analysis'

            # remediation_summary — prefer the explicit DB field, then fall through
            # to every other field that holds the AI's reasoning for the gap.
            remediation_summary = (
                form.get('remediation_summary')
                or form.get('remediation')
                or (ai_results.get('recommendation') if isinstance(ai_results, dict) else None)
                or form.get('implementation_description')
                or ("; ".join(gaps) if gaps else None)
                or ''
            )

            last_updated = form.get('last_modified_at') or form.get('submitted_at')
            if hasattr(last_updated, 'isoformat'):
                last_updated = last_updated.isoformat()

            findings_results.append({
                # ── Identity ──────────────────────────────────────────────────
                'id':                  form.get('id'),
                'submission_id':       form.get('id'),
                'evidence_id':         None,
                'source_type':         'intake_form',
                # ── Control metadata ──────────────────────────────────────────
                'control_id':          form.get('control_id'),
                'control':             form.get('control_id'),
                'control_name':        form.get('control_title') or form.get('control_id'),
                'framework':           form.get('framework_id') or 'Unknown',
                # ── Three critical PDF fields ─────────────────────────────────
                'status':              status_val,
                'remediation_summary': remediation_summary,
                'ai_confidence_score': confidence,              # raw 0.0–1.0
                'aiConfidence':        int(confidence * 100),   # percent int for PDF
                # ── Supporting detail ─────────────────────────────────────────
                'description':         form.get('implementation_description') or remediation_summary,
                'gap_summary':         "; ".join(gaps) if isinstance(gaps, list) else str(gaps),
                'recommendation':      (ai_results.get('recommendation') if isinstance(ai_results, dict) else None) or remediation_summary,
                'severity':            (ai_results.get('risk_level', 'Medium').capitalize() if isinstance(ai_results, dict) else 'Medium'),
                'analyst_id':          form.get('analyst_id'),
                'analyst_name':        analyst_map.get(form.get('analyst_id'), 'Unknown Analyst'),
                'analyst_description': form.get('implementation_description'),
                'auditor_comments':    form.get('auditor_comments'),
                'auditor_name':        form.get('auditor_name'),
                'ai_details':          ai_results,
                'form_responses':      form.get('form_data') or form.get('responses') or {},
                'last_updated':        last_updated,
            })

        # 2. mvp1_evidence → findings (secondary source)
        for ev in evidence_list:
            if not isinstance(ev, dict):
                continue

            raw_status = ev.get('status') or 'uploaded'
            confidence = float(ev.get('ai_confidence_score') or 0.0)
            ai_res     = ev.get('ai_analysis_results') or {}

            EVIDENCE_TERMINAL = {'approved', 'rejected', 'analyzed', 'requires_revision'}
            if raw_status in EVIDENCE_TERMINAL:
                status_val = 'ai_verified' if raw_status == 'analyzed' else raw_status
            elif confidence >= 0.8 and (ai_res.get('risk_level') or '').lower() != 'high':
                status_val = 'ai_verified'
            elif confidence > 0:
                status_val = 'pending'
            else:
                status_val = 'pending_ai_analysis'

            remediation_summary = (
                ev.get('remediation_summary')
                or (ai_res.get('recommendation') if isinstance(ai_res, dict) else None)
                or ev.get('description')
                or ev.get('notes')
                or ''
            )

            ev_last = ev.get('uploaded_at') or ev.get('updated_at')
            if hasattr(ev_last, 'isoformat'):
                ev_last = ev_last.isoformat()

            findings_results.append({
                'id':                  ev.get('id'),
                'submission_id':       ev.get('id'),
                'evidence_id':         ev.get('id'),
                'source_type':         'evidence',
                'control_id':          ev.get('control_id'),
                'control':             ev.get('control_id'),
                'control_name':        ev.get('control_name') or ev.get('control_id'),
                'filename':            ev.get('original_filename'),
                'framework':           'Evidence Analysis',
                # ── Three critical PDF fields ─────────────────────────────────
                'status':              status_val,
                'remediation_summary': remediation_summary,
                'ai_confidence_score': confidence,
                'aiConfidence':        int(confidence * 100),
                # ── Supporting detail ─────────────────────────────────────────
                'description':         ev.get('description') or remediation_summary,
                'gap_summary':         '',
                'recommendation':      remediation_summary,
                'severity':            (ai_res.get('risk_level', 'Medium').capitalize() if isinstance(ai_res, dict) else 'Medium'),
                'analyst_description': ev.get('description'),
                'ai_details':          ai_res,
                'last_updated':        ev_last,
            })

        # 3. Legacy control_assessments → findings (fallback)
        for ctrl in (data.get('control_assessments') or []):
            if not isinstance(ctrl, dict):
                continue
            confidence = float(ctrl.get('ai_confidence') or 0.0)
            remediation_summary = (
                ctrl.get('remediation_summary')
                or ctrl.get('recommendation')
                or ctrl.get('gap_summary')
                or ''
            )
            findings_results.append({
                'id':                  ctrl.get('id', str(uuid.uuid4())),
                'source_type':         'finding',
                'control_id':          ctrl.get('control_id'),
                'control':             ctrl.get('control_name') or ctrl.get('control_id'),
                'control_name':        ctrl.get('control_name') or '',
                'framework':           ctrl.get('framework') or 'Unknown',
                'status':              ctrl.get('status') or 'pending',
                'remediation_summary': remediation_summary,
                'ai_confidence_score': confidence,
                'aiConfidence':        int(confidence * 100),
                'description':         ctrl.get('description') or '',
                'gap_summary':         ctrl.get('gap_summary') or '',
                'recommendation':      ctrl.get('recommendation') or 'Review required',
                'severity':            'Medium',
                'last_updated':        ctrl.get('updated_at'),
            })

        data['generated_findings'] = findings_results
        logger.info(
            f"[report {config.id}] generated_findings total: {len(findings_results)} "
            f"({len(intake_forms)} intake forms + {len(evidence_list)} evidence items + "
            f"{len(data.get('control_assessments', []))} legacy controls)"
        )

        # ── Optional supplementary collections ───────────────────────────────────
        for coll, key, query_field in [
            ("mvp1_auditor_feedback",    "mvp1_auditor_feedback",    "engagement_id"),
            ("mvp1_evidence_annotations","mvp1_evidence_annotations", "engagement_id"),
            ("gap_validations",          "gap_validations",           "engagement_id"),
            ("ai_gap_findings",          "ai_gap_findings",           "assessment_id"),
        ]:
            try:
                qid = config.assessment_id
                rows = await DatabaseOperations.find_many(coll, {query_field: qid})
                data[key] = rows or []
            except Exception as misc_err:
                logger.debug(f"Collection '{coll}' not available: {misc_err}")
                data[key] = []

        # ── Report creator ───────────────────────────────────────────────────────
        try:
            if config.created_by:
                user = await DatabaseOperations.find_one("users", {"id": config.created_by})
                if user:
                    data['created_by'] = user
        except Exception:
            pass

        return data




    async def _generate_tokuro_insights(self, report_data: Dict[str, Any]) -> TokuroAIReportData:
        """Generate AI-powered insights for the report"""
        
        try:
            # Prepare assessment data for AI analysis
            assessment = report_data.get('assessment', {})
            controls = report_data.get('control_assessments', [])
            evidence = report_data.get('evidence_files', [])
            
            # Generate AI insights
            insights_prompt = f"""
            Generate comprehensive compliance report insights for:
            Assessment: {assessment.get('title', 'N/A')}
            Framework: {assessment.get('framework', 'NIST CSF 2.0')}
            Controls Assessed: {len(controls)}
            Evidence Files: {len(evidence)}
            
            Provide executive-level insights including:
            1. Overall compliance posture
            2. Key risk areas
            3. Remediation priorities
            4. Strategic recommendations
            """
            
            ai_response = await self.tokuro_ai.analyze_text(insights_prompt)
            
            # Calculate confidence scores
            confidence_scores = {}
            for control in controls:
                if control.get('ai_confidence'):
                    confidence_scores[control.get('control_id', 'unknown')] = control['ai_confidence']
            
            # Extract framework mappings
            framework_mappings = {}
            for evidence_file in evidence:
                if evidence_file.get('framework_mappings'):
                    filename = evidence_file.get('filename', 'unknown')
                    framework_mappings[filename] = evidence_file['framework_mappings']
            
            # Calculate maturity scores
            maturity_scores = {}
            for control in controls:
                control_id = control.get('control_id', 'unknown')
                maturity_scores[control_id] = control.get('maturity_score', 0) / 5.0  # Normalize to 0-1
            
            # Generate remediation recommendations
            remediation_recommendations = []
            for control in controls:
                if control.get('overall_score', 0) < 0.7:  # Low scoring controls
                    remediation_recommendations.append({
                        'control_id': control.get('control_id'),
                        'priority': 'High' if control.get('overall_score', 0) < 0.5 else 'Medium',
                        'recommendation': f"Improve implementation of {control.get('control_id')}",
                        'estimated_effort': 'Medium'
                    })
            
            return TokuroAIReportData(
                confidence_scores=confidence_scores,
                framework_mappings=framework_mappings,
                gap_analysis={
                    'total_gaps': len(remediation_recommendations),
                    'critical_gaps': len([r for r in remediation_recommendations if r['priority'] == 'High']),
                    'summary': ai_response if isinstance(ai_response, str) else 'AI analysis completed'
                },
                remediation_recommendations=remediation_recommendations,
                maturity_scores=maturity_scores,
                risk_assessment={
                    'overall_risk_level': 'Medium',
                    'risk_areas': ['Access Control', 'Data Protection', 'Incident Response'],
                    'risk_score': 6.5
                },
                executive_insights=[
                    'Organization demonstrates strong commitment to cybersecurity compliance',
                    'Key improvement areas identified in access control and data protection',
                    'Recommend prioritizing high-risk remediation activities',
                    'Overall compliance trajectory is positive with targeted improvements needed'
                ]
            )
            
        except Exception as e:
            logger.error(f"Error generating Tokuro AI insights: {str(e)}")
            # Return minimal insights on error
            return TokuroAIReportData()

    async def _load_template(self, template_id: str) -> ReportTemplate:
        """Load report template configuration"""
        
        try:
            # Try to load from database first
            template = await DatabaseOperations.find_one(
                "report_templates",
                {"id": template_id}
            )
            
            if template:
                return ReportTemplate(**template)
            
            # Fall back to default templates
            return self._get_default_template(template_id)
            
        except Exception as e:
            logger.error(f"Error loading template {template_id}: {str(e)}")
            # Return basic template
            return self._get_default_template("basic")

    def _get_default_template(self, template_type: str = "basic") -> ReportTemplate:
        """Get default report template"""
        
        templates = {
            "basic": ReportTemplate(
                id="default_basic",
                name="Basic Assessment Report",
                description="Standard assessment report with key findings",
                report_type=ReportType.ASSESSMENT_SUMMARY,
                supported_formats=[ReportFormat.PDF, ReportFormat.PPTX],
                template_path="basic_template.html",
                sections=['executive_summary', 'assessment_overview', 'key_findings', 'recommendations'],
                dynamic_fields={
                    'organization_name': '{{organization.name}}',
                    'assessment_title': '{{assessment.title}}',
                    'generated_date': '{{report_info.generated_at}}',
                    'total_controls': '{{control_assessments|length}}'
                }
            ),
            "executive": ReportTemplate(
                id="default_executive",
                name="Executive Summary Report",
                description="High-level executive presentation",
                report_type=ReportType.EXECUTIVE_SUMMARY,
                supported_formats=[ReportFormat.PPTX, ReportFormat.PDF],
                template_path="executive_template.html",
                sections=['executive_summary', 'risk_overview', 'strategic_recommendations'],
                dynamic_fields={
                    'organization_name': '{{organization.name}}',
                    'compliance_score': '{{assessment.compliance_score}}',
                    'risk_level': '{{tokuro_ai_data.risk_assessment.overall_risk_level}}'
                }
            )
        }
        
        return templates.get(template_type, templates["basic"])

    async def _generate_by_format(
        self, 
        config: ReportGeneration, 
        data: Dict[str, Any], 
        template: ReportTemplate
    ) -> str:
        """Generate report based on specified format"""
        
        format_generators = {
            ReportFormat.PDF: self._generate_pdf_report,
            ReportFormat.PPTX: self._generate_pptx_report,
            ReportFormat.DOCX: self._generate_docx_report,
            ReportFormat.XLSX: self._generate_xlsx_report
        }
        
        generator = format_generators.get(config.report_format)
        if not generator:
            raise ValueError(f"Unsupported report format: {config.report_format}")
        
        return await generator(config, data, template)

    async def _generate_pdf_report(
        self, 
        config: ReportGeneration, 
        data: Dict[str, Any], 
        template: ReportTemplate
    ) -> str:
        """Generate PDF report using ReportLab"""
        
        filename = f"{config.id}_{config.name.replace(' ', '_')}.pdf"
        output_path = self.output_dir / filename
        
        # Create PDF document
        doc = SimpleDocTemplate(
            str(output_path),
            pagesize=A4,
            topMargin=1*inch,
            bottomMargin=1*inch,
            leftMargin=1*inch,
            rightMargin=1*inch
        )
        
        # Build content
        story = []
        styles = getSampleStyleSheet()
        
        # Custom styles
        title_style = ParagraphStyle(
            'CustomTitle',
            parent=styles['Title'],
            fontSize=24,
            textColor=colors.HexColor(self.branding['primary_color']),
            spaceAfter=30
        )
        
        heading_style = ParagraphStyle(
            'CustomHeading',
            parent=styles['Heading1'],
            fontSize=16,
            textColor=colors.HexColor(self.branding['secondary_color']),
            spaceAfter=12
        )
        
        # Title page
        story.append(Paragraph(data['report_info']['name'], title_style))
        story.append(Spacer(1, 20))
        
        # Organization info
        if data.get('organization'):
            story.append(Paragraph(f"Organization: {data['organization'].get('name', 'N/A')}", styles['Normal']))
            story.append(Spacer(1, 12))
        
        # Generated date and Tokuro AI branding
        story.append(Paragraph(f"Generated on: {data['report_info']['generated_at'].strftime('%B %d, %Y')}", styles['Normal']))
        story.append(Paragraph(f"<i>{self.branding['watermark_text']}</i>", styles['Italic']))
        story.append(Spacer(1, 30))
        
        # Executive Summary
        story.append(Paragraph("Executive Summary", heading_style))
        if config.tokuro_ai_data and config.tokuro_ai_data.executive_insights:
            for insight in config.tokuro_ai_data.executive_insights:
                story.append(Paragraph(f"• {insight}", styles['Normal']))
                story.append(Spacer(1, 6))
        else:
            story.append(Paragraph("Executive summary will be populated with assessment data.", styles['Normal']))
        story.append(Spacer(1, 20))
        
        # Assessment Overview
        story.append(Paragraph("Assessment Overview", heading_style))
        if data.get('assessment'):
            assessment = data['assessment']
            story.append(Paragraph(f"<b>Title:</b> {assessment.get('title', 'N/A')}", styles['Normal']))
            story.append(Paragraph(f"<b>Framework:</b> {assessment.get('framework', 'N/A')}", styles['Normal']))
            story.append(Paragraph(f"<b>Status:</b> {assessment.get('status', 'N/A')}", styles['Normal']))
            story.append(Spacer(1, 6))
        
        # Control Assessment Summary
        controls = data.get('control_assessments', [])
        if controls:
            story.append(Paragraph("Control Assessment Summary", heading_style))
            
            # Create summary table
            table_data = [['Control ID', 'Status', 'Score', 'Maturity']]
            for control in controls[:10]:  # Limit to first 10 for space
                table_data.append([
                    control.get('control_id', 'N/A'),
                    control.get('status', 'N/A'),
                    f"{control.get('overall_score', 0):.1f}",
                    f"{control.get('maturity_score', 0)}/5"
                ])
            
            table = Table(table_data)
            table.setStyle(TableStyle([
                ('BACKGROUND', (0, 0), (-1, 0), colors.HexColor(self.branding['secondary_color'])),
                ('TEXTCOLOR', (0, 0), (-1, 0), colors.whitesmoke),
                ('ALIGN', (0, 0), (-1, -1), 'LEFT'),
                ('FONTNAME', (0, 0), (-1, 0), 'Helvetica-Bold'),
                ('FONTSIZE', (0, 0), (-1, 0), 12),
                ('BOTTOMPADDING', (0, 0), (-1, 0), 12),
                ('BACKGROUND', (0, 1), (-1, -1), colors.beige),
                ('GRID', (0, 0), (-1, -1), 1, colors.black)
            ]))
            
            story.append(table)
            story.append(Spacer(1, 20))

        # -----------------------------------------------------------------------
        # REVIEW STATE: Summary + Findings Tables (always rendered unconditionally)
        # -----------------------------------------------------------------------
        findings = data.get('generated_findings') or []

        # Define cell/body paragraph style for wrapping text in table cells
        cell_style = ParagraphStyle(
            'TableCell',
            parent=styles['Normal'],
            fontSize=9,
            leading=12,
            wordWrap='LTR',
        )
        header_style = ParagraphStyle(
            'TableHeader',
            parent=styles['Normal'],
            fontSize=10,
            leading=13,
            textColor=colors.whitesmoke,
            fontName='Helvetica-Bold',
        )

        # ---
        # 1. SUMMARY TABLE
        # ---
        total_controls = len({f.get('control_id') for f in findings if f.get('control_id')})
        total_gaps = len([
            f for f in findings
            if (f.get('status') or '').lower() in ('gap_found', 'flagged', 'needs_remediation')
        ])
        total_pending = len([
            f for f in findings
            if (f.get('status') or '').lower() in ('pending', 'requires_revision', 'pending_ai_analysis')
        ])

        story.append(Paragraph('Review State Summary', heading_style))
        summary_tbl_data = [
            [Paragraph('Metric', header_style),      Paragraph('Count', header_style)],
            [Paragraph('Total Controls', cell_style), Paragraph(str(total_controls), cell_style)],
            [Paragraph('Total Gaps Identified', cell_style), Paragraph(str(total_gaps), cell_style)],
            [Paragraph('Items Pending Review', cell_style), Paragraph(str(total_pending), cell_style)],
        ]
        summary_tbl = Table(summary_tbl_data, colWidths=[3.5 * inch, 2 * inch])
        summary_tbl.setStyle(TableStyle([
            ('BACKGROUND',    (0, 0), (-1, 0), colors.HexColor(self.branding['primary_color'])),
            ('TEXTCOLOR',     (0, 0), (-1, 0), colors.whitesmoke),
            ('FONTNAME',      (0, 0), (-1, 0), 'Helvetica-Bold'),
            ('FONTSIZE',      (0, 0), (-1, 0), 11),
            ('GRID',          (0, 0), (-1, -1), 0.5, colors.grey),
            ('LEFTPADDING',   (0, 0), (-1, -1), 8),
            ('RIGHTPADDING',  (0, 0), (-1, -1), 8),
            ('TOPPADDING',    (0, 0), (-1, -1), 5),
            ('BOTTOMPADDING', (0, 0), (-1, -1), 5),
            ('VALIGN',        (0, 0), (-1, -1), 'MIDDLE'),
        ]))
        story.append(summary_tbl)
        story.append(Spacer(1, 20))

        # ---
        # 2. IDENTIFIED GAPS TABLE (status: gap_found | flagged | needs_remediation)
        # ---
        story.append(Paragraph('Identified Gaps', heading_style))
        gap_rows = [
            [
                Paragraph('Control ID', header_style),
                Paragraph('AI Reason / Description', header_style),
                Paragraph('Remediation Summary', header_style),
            ]
        ]
        for f in findings:
            st = (f.get('status') or '').lower()
            if st in ('gap_found', 'flagged', 'needs_remediation'):
                cid         = str(f.get('control_id') or f.get('control') or 'N/A')
                ai_reason   = str(
                    f.get('description')
                    or (f.get('ai_details') or {}).get('summary')
                    or (f.get('ai_details') or {}).get('gap_reason')
                    or '—'
                )
                remediation = str(
                    f.get('remediation_summary')
                    or f.get('remediation')
                    or f.get('recommendation')
                    or f.get('auditor_comments')
                    or '—'
                )
                gap_rows.append([
                    Paragraph(cid,         cell_style),
                    Paragraph(ai_reason,   cell_style),
                    Paragraph(remediation, cell_style),
                ])

        if len(gap_rows) == 1:
            story.append(Paragraph('No identified gaps for this engagement.', styles['Normal']))
        else:
            gap_tbl = Table(gap_rows, colWidths=[1.4 * inch, 3.0 * inch, 2.2 * inch])
            gap_tbl.setStyle(TableStyle([
                ('BACKGROUND',    (0, 0), (-1, 0), colors.HexColor(self.branding['primary_color'])),
                ('TEXTCOLOR',     (0, 0), (-1, 0), colors.whitesmoke),
                ('FONTNAME',      (0, 0), (-1, 0), 'Helvetica-Bold'),
                ('FONTSIZE',      (0, 0), (-1, 0), 10),
                ('GRID',          (0, 0), (-1, -1), 0.5, colors.grey),
                ('LEFTPADDING',   (0, 0), (-1, -1), 6),
                ('RIGHTPADDING',  (0, 0), (-1, -1), 6),
                ('TOPPADDING',    (0, 0), (-1, -1), 5),
                ('BOTTOMPADDING', (0, 0), (-1, -1), 5),
                ('VALIGN',        (0, 0), (-1, -1), 'TOP'),
                ('ROWBACKGROUNDS', (0, 1), (-1, -1), [colors.white, colors.HexColor('#f3f4f6')]),
            ]))
            story.append(gap_tbl)

        story.append(Spacer(1, 18))

        # ---
        # 3. AI VERIFIED & APPROVED TABLE (status: ai_verified | approved)
        # ---
        story.append(Paragraph('AI Verified & Approved Controls', heading_style))
        verified_rows = [
            [
                Paragraph('Control ID', header_style),
                Paragraph('Control Name', header_style),
                Paragraph('Status', header_style),
                Paragraph('AI Confidence', header_style),
            ]
        ]
        for f in findings:
            st = (f.get('status') or '').lower()
            if st in ('ai_verified', 'approved'):
                try:
                    ai_conf_raw = f.get('aiConfidence')
                    if ai_conf_raw is not None:
                        ai_conf_pct = f"{int(ai_conf_raw)}%"
                    else:
                        ai_conf_pct = f"{int((f.get('ai_confidence_score') or 0) * 100)}%"
                except Exception:
                    ai_conf_pct = '0%'

                cid   = str(f.get('control_id') or f.get('control') or 'N/A')
                cname = str(f.get('control_name') or f.get('control') or '—')
                verified_rows.append([
                    Paragraph(cid,                    cell_style),
                    Paragraph(cname,                  cell_style),
                    Paragraph(str(f.get('status', '')), cell_style),
                    Paragraph(ai_conf_pct,            cell_style),
                ])

        if len(verified_rows) == 1:
            story.append(Paragraph('No AI-verified or approved controls for this engagement.', styles['Normal']))
        else:
            ver_tbl = Table(verified_rows, colWidths=[1.4 * inch, 2.8 * inch, 1.2 * inch, 1.2 * inch])
            ver_tbl.setStyle(TableStyle([
                ('BACKGROUND',    (0, 0), (-1, 0), colors.HexColor(self.branding['primary_color'])),
                ('TEXTCOLOR',     (0, 0), (-1, 0), colors.whitesmoke),
                ('FONTNAME',      (0, 0), (-1, 0), 'Helvetica-Bold'),
                ('FONTSIZE',      (0, 0), (-1, 0), 10),
                ('GRID',          (0, 0), (-1, -1), 0.5, colors.grey),
                ('LEFTPADDING',   (0, 0), (-1, -1), 6),
                ('RIGHTPADDING',  (0, 0), (-1, -1), 6),
                ('TOPPADDING',    (0, 0), (-1, -1), 5),
                ('BOTTOMPADDING', (0, 0), (-1, -1), 5),
                ('VALIGN',        (0, 0), (-1, -1), 'TOP'),
                ('ROWBACKGROUNDS', (0, 1), (-1, -1), [colors.white, colors.HexColor('#f3f4f6')]),
            ]))
            story.append(ver_tbl)

        story.append(Spacer(1, 24))

        # -----------------------------------------------------------------------
        # 4. DETAILED FINDINGS — full per-finding audit table (mirrors ReviewState.jsx)
        # -----------------------------------------------------------------------
        # Status → human-readable label map (synced with ReviewState.jsx logic)
        STATUS_LABEL_MAP = {
            'gap_found':          'Identified Gap',
            'flagged':            'Identified Gap',
            'needs_remediation':  'Identified Gap',
            'ai_verified':        'AI Verified',
            'approved':           'Auditor Approved',
            'rejected':           'Rejected',
            'pending':            'Pending Review',
            'requires_revision':  'Requires Revision',
            'pending_ai_analysis':'Awaiting AI Analysis',
        }

        # Severity badge colours used inside status cells
        STATUS_COLORS = {
            'Identified Gap':       colors.HexColor('#dc2626'),   # red
            'Rejected':             colors.HexColor('#7f1d1d'),   # dark red
            'Pending Review':       colors.HexColor('#d97706'),   # amber
            'Requires Revision':    colors.HexColor('#b45309'),   # dark amber
            'Awaiting AI Analysis': colors.HexColor('#6b7280'),   # grey
            'AI Verified':          colors.HexColor('#059669'),   # green
            'Auditor Approved':     colors.HexColor('#1d4ed8'),   # blue
        }

        story.append(Paragraph('Detailed Findings', heading_style))

        if not findings:
            story.append(Paragraph(
                'No findings have been generated for this engagement yet.',
                styles['Normal']
            ))
        else:
            # --- Define a small style for the status badge text ---
            badge_normal = ParagraphStyle(
                'BadgeNormal',
                parent=styles['Normal'],
                fontSize=8,
                leading=11,
                textColor=colors.white,
                fontName='Helvetica-Bold',
                wordWrap='LTR',
            )
            detail_label = ParagraphStyle(
                'DetailLabel',
                parent=styles['Normal'],
                fontSize=8,
                leading=11,
                textColor=colors.HexColor('#374151'),
                fontName='Helvetica-Bold',
            )
            detail_body = ParagraphStyle(
                'DetailBody',
                parent=styles['Normal'],
                fontSize=8,
                leading=11,
                textColor=colors.HexColor('#111827'),
                wordWrap='LTR',
            )

            # Column widths for A4 with 1" margins = 6.27" usable.
            # Total here = 1.0 + 1.5 + 1.1 + 0.8 + 1.87 = 6.27"
            COL_WIDTHS = [1.0 * inch, 1.5 * inch, 1.1 * inch, 0.8 * inch, 1.87 * inch]

            detail_rows = [
                [
                    Paragraph('Control ID',          header_style),
                    Paragraph('Control Name',         header_style),
                    Paragraph('Review Status',        header_style),
                    Paragraph('AI Confidence',        header_style),
                    Paragraph('Remediation Summary',  header_style),
                ]
            ]

            # Track per-row background overrides for status colouring
            row_bg_overrides = []   # list of (row_index, hex_color_string)

            for finding in findings:
                raw_status  = (finding.get('status') or '').lower()
                label       = STATUS_LABEL_MAP.get(raw_status, raw_status.replace('_', ' ').title())

                # ── AI Confidence ──────────────────────────────────────────────
                try:
                    ai_conf_raw = finding.get('aiConfidence')
                    if ai_conf_raw is not None:
                        ai_conf_pct = f"{int(ai_conf_raw)}%"
                    else:
                        score = finding.get('ai_confidence_score') or 0
                        ai_conf_pct = f"{int(float(score) * 100)}%"
                except Exception:
                    ai_conf_pct = '0%'

                # ── Field extraction ───────────────────────────────────────────
                control_id   = str(finding.get('control_id') or finding.get('control') or '—')
                control_name = str(
                    finding.get('control_name')
                    or finding.get('control_title')
                    or finding.get('control')
                    or '—'
                )

                # remediation_summary is the primary field; fall through to every
                # plausible alias so nothing is ever left blank for an auditor.
                remediation_text = str(
                    finding.get('remediation_summary')
                    or finding.get('remediation')
                    or finding.get('recommendation')
                    or finding.get('description')
                    or finding.get('gap_summary')
                    or finding.get('auditor_comments')
                    or (finding.get('ai_details') or {}).get('recommendation')
                    or (finding.get('ai_details') or {}).get('summary')
                    or '—'
                )

                row_data = [
                    Paragraph(control_id,      cell_style),
                    Paragraph(control_name,    cell_style),
                    Paragraph(label,           cell_style),
                    Paragraph(ai_conf_pct,     cell_style),
                    Paragraph(remediation_text, cell_style),
                ]
                detail_rows.append(row_data)

            detail_tbl = Table(detail_rows, colWidths=COL_WIDTHS, repeatRows=1)
            detail_tbl.setStyle(TableStyle([
                # Header row
                ('BACKGROUND',    (0, 0), (-1, 0),  colors.HexColor(self.branding['primary_color'])),
                ('TEXTCOLOR',     (0, 0), (-1, 0),  colors.whitesmoke),
                ('FONTNAME',      (0, 0), (-1, 0),  'Helvetica-Bold'),
                ('FONTSIZE',      (0, 0), (-1, 0),  10),
                # Data rows
                ('FONTNAME',      (0, 1), (-1, -1), 'Helvetica'),
                ('FONTSIZE',      (0, 1), (-1, -1), 8),
                # Alternating row shading
                ('ROWBACKGROUNDS',(0, 1), (-1, -1), [colors.white, colors.HexColor('#f9fafb')]),
                # Grid and padding
                ('GRID',          (0, 0), (-1, -1), 0.4, colors.HexColor('#d1d5db')),
                ('LEFTPADDING',   (0, 0), (-1, -1), 6),
                ('RIGHTPADDING',  (0, 0), (-1, -1), 6),
                ('TOPPADDING',    (0, 0), (-1, -1), 5),
                ('BOTTOMPADDING', (0, 0), (-1, -1), 5),
                # Critical: top-align so multiline remediation stays next to its Control ID
                ('VALIGN',        (0, 0), (-1, -1), 'TOP'),
                # Highlight the Status column header slightly
                ('ALIGN',         (3, 0), (3, -1),  'CENTER'),   # AI Confidence centred
            ]))

            story.append(detail_tbl)

        story.append(Spacer(1, 24))

        # -----------------------------------------------------------------------
        # 5. DETAILED COMPLIANCE REVIEW — compact 4-column findings table
        # -----------------------------------------------------------------------
        story.append(Paragraph("Detailed Compliance Review", heading_style))

        # Table header
        findings_table_data = [["Control ID", "Review Status", "Confidence", "Detailed Findings & Gaps"]]

        for f in data.get('generated_findings', []):
            control = Paragraph(
                f"<b>{f.get('control_id', 'N/A')}</b><br/>{f.get('control_name', '')}",
                styles['Normal']
            )
            status = (f.get('status') or '').replace('_', ' ').title()
            conf   = f"{f.get('aiConfidence') or int((f.get('ai_confidence_score') or 0) * 100)}%"

            # remediation_summary is the critical data point — fall back through
            # every alias so the auditor always sees the AI's reasoning for the gap.
            findings_text = (
                f.get('remediation_summary')
                or f.get('description')
                or f.get('recommendation')
                or f.get('gap_summary')
                or 'N/A'
            )
            findings_box = Paragraph(findings_text, styles['Normal'])

            findings_table_data.append([control, status, conf, findings_box])

        # col widths: 1.1 + 1.2 + 0.8 + 3.1 = 6.2" — fits A4 with 1" margins
        findings_table = Table(
            findings_table_data,
            colWidths=[1.1 * inch, 1.2 * inch, 0.8 * inch, 3.1 * inch],
            repeatRows=1,
        )
        findings_table.setStyle(TableStyle([
            ('BACKGROUND',    (0, 0), (-1, 0), colors.HexColor(self.branding['primary_color'])),
            ('TEXTCOLOR',     (0, 0), (-1, 0), colors.white),
            ('FONTNAME',      (0, 0), (-1, 0), 'Helvetica-Bold'),
            ('FONTSIZE',      (0, 0), (-1, 0), 10),
            ('VALIGN',        (0, 0), (-1, -1), 'TOP'),
            ('GRID',          (0, 0), (-1, -1), 0.5, colors.grey),
            ('LEFTPADDING',   (0, 0), (-1, -1), 6),
            ('RIGHTPADDING',  (0, 0), (-1, -1), 6),
            ('TOPPADDING',    (0, 0), (-1, -1), 5),
            ('BOTTOMPADDING', (0, 0), (-1, -1), 5),
            ('ROWBACKGROUNDS',(0, 1), (-1, -1), [colors.white, colors.HexColor('#f3f4f6')]),
        ]))

        story.append(findings_table)
        story.append(Spacer(1, 24))

        # -----------------------------------------------------------------------
        # TOKURO AI INSIGHTS (optional — only if AI data was generated)
        # -----------------------------------------------------------------------
        if config.tokuro_ai_data:
            story.append(Paragraph("Tokuro AI Insights", heading_style))

            tokuro_data = config.tokuro_ai_data
            if tokuro_data.gap_analysis:
                gap = tokuro_data.gap_analysis
                story.append(Paragraph(f"<b>Total Gaps Identified:</b> {gap.get('total_gaps', 0)}", styles['Normal']))
                story.append(Paragraph(f"<b>Critical Gaps:</b> {gap.get('critical_gaps', 0)}", styles['Normal']))
                story.append(Spacer(1, 12))

            if tokuro_data.remediation_recommendations:
                story.append(Paragraph("<b>Top Recommendations:</b>", styles['Normal']))
                for i, rec in enumerate(tokuro_data.remediation_recommendations[:5]):
                    story.append(Paragraph(
                        f"{i + 1}. {rec.get('recommendation', 'N/A')} (Priority: {rec.get('priority', 'N/A')})",
                        styles['Normal']
                    ))
                    story.append(Spacer(1, 6))

            story.append(Spacer(1, 18))

        # -----------------------------------------------------------------------
        # INTAKE FORMS — full submitted responses
        # -----------------------------------------------------------------------
        story.append(Paragraph("Intake Forms & Submitted Evidence", heading_style))
        intake_forms = data.get('intake_forms', [])
        if intake_forms:
            for form in intake_forms:
                form_title = form.get('title', form.get('id', 'Intake Form'))
                story.append(Paragraph(f"<b>Form:</b> {form_title}", styles['Normal']))
                responses = form.get('responses') or form.get('form_responses') or form
                try:
                    if isinstance(responses, dict):
                        for k, v in responses.items():
                            story.append(Paragraph(f"&nbsp;&nbsp;<b>{k}:</b> {str(v)}", styles['Normal']))
                    else:
                        story.append(Paragraph(str(responses), styles['Normal']))
                except Exception:
                    story.append(Paragraph("(Could not render form responses)", styles['Normal']))
                story.append(Spacer(1, 6))
        else:
            story.append(Paragraph("No intake forms submitted for this engagement.", styles['Normal']))

        # -----------------------------------------------------------------------
        # EVIDENCE FILES — metadata listing
        # -----------------------------------------------------------------------
        story.append(Spacer(1, 12))
        story.append(Paragraph("Submitted Evidence Files", heading_style))
        evidence_files = data.get('evidence_files', [])
        if evidence_files:
            for ef in evidence_files:
                fname    = ef.get('filename') or ef.get('name') or ef.get('id') or 'Unknown'
                uploader = ef.get('uploaded_by') or ef.get('uploader') or ef.get('created_by') or 'Unknown'
                story.append(Paragraph(f"- <b>{fname}</b> (Uploaded by: {uploader})", styles['Normal']))
                if ef.get('description'):
                    story.append(Paragraph(f"&nbsp;&nbsp;{ef.get('description')}", styles['Normal']))
                # Embed small images if available on disk
                file_path = ef.get('file_path') or ef.get('path') or ef.get('local_path')
                try:
                    if file_path and os.path.exists(file_path) and str(fname).lower().endswith(('.png', '.jpg', '.jpeg')):
                        story.append(Image(file_path, width=4 * inch, height=3 * inch))
                except Exception:
                    pass
                story.append(Spacer(1, 6))
        else:
            story.append(Paragraph("No evidence files attached to this engagement.", styles['Normal']))

        # Build PDF
        doc.build(story)

        return str(output_path)

    async def _generate_pptx_report(
        self, 
        config: ReportGeneration, 
        data: Dict[str, Any], 
        template: ReportTemplate
    ) -> str:
        """Generate PowerPoint presentation"""
        
        filename = f"{config.id}_{config.name.replace(' ', '_')}.pptx"
        output_path = self.output_dir / filename
        
        # Create presentation
        prs = Presentation()
        
        # Set slide dimensions (16:9 aspect ratio)
        prs.slide_width = Inches(13.33)
        prs.slide_height = Inches(7.5)
        
        # Title slide
        slide_layout = prs.slide_layouts[0]  # Title slide
        slide = prs.slides.add_slide(slide_layout)
        
        # Configure title
        title = slide.shapes.title
        title.text = data['report_info']['name']
        title.text_frame.paragraphs[0].font.size = Pt(36)
        title.text_frame.paragraphs[0].font.color.rgb = RGBColor.from_string(self.branding['primary_color'].replace('#', ''))
        
        # Subtitle
        subtitle = slide.placeholders[1]
        subtitle_text = f"{data['organization'].get('name', 'Organization')}\n"
        subtitle_text += f"Generated on {data['report_info']['generated_at'].strftime('%B %d, %Y')}\n"
        subtitle_text += f"{self.branding['watermark_text']}"
        subtitle.text = subtitle_text
        
        # Executive Summary slide
        slide_layout = prs.slide_layouts[1]  # Title and content
        slide = prs.slides.add_slide(slide_layout)
        title = slide.shapes.title
        title.text = "Executive Summary"
        
        content = slide.placeholders[1].text_frame
        content.clear()
        
        if config.tokuro_ai_data and config.tokuro_ai_data.executive_insights:
            for insight in config.tokuro_ai_data.executive_insights:
                p = content.add_paragraph()
                p.text = insight
                p.level = 0
        else:
            p = content.add_paragraph()
            p.text = "Executive insights will be generated based on assessment data"
            p.level = 0
        
        # Assessment Overview slide
        slide = prs.slides.add_slide(slide_layout)
        title = slide.shapes.title
        title.text = "Assessment Overview"
        
        content = slide.placeholders[1].text_frame
        content.clear()
        
        if data.get('assessment'):
            assessment = data['assessment']
            
            p = content.add_paragraph()
            p.text = f"Assessment: {assessment.get('title', 'N/A')}"
            p.level = 0
            
            p = content.add_paragraph()
            p.text = f"Framework: {assessment.get('framework', 'N/A')}"
            p.level = 0
            
            p = content.add_paragraph()
            p.text = f"Status: {assessment.get('status', 'N/A')}"
            p.level = 0
            
            p = content.add_paragraph()
            p.text = f"Controls Assessed: {len(data.get('control_assessments', []))}"
            p.level = 0
            
            p = content.add_paragraph()
            p.text = f"Evidence Files: {len(data.get('evidence_files', []))}"
            p.level = 0
        
        # Key Findings slide
        if config.tokuro_ai_data:
            slide = prs.slides.add_slide(slide_layout)
            title = slide.shapes.title
            title.text = "Key Findings & Recommendations"
            
            content = slide.placeholders[1].text_frame
            content.clear()
            
            tokuro_data = config.tokuro_ai_data
            
            if tokuro_data.gap_analysis:
                gap = tokuro_data.gap_analysis
                p = content.add_paragraph()
                p.text = f"Total Gaps Identified: {gap.get('total_gaps', 0)}"
                p.level = 0
                
                p = content.add_paragraph()
                p.text = f"Critical Gaps: {gap.get('critical_gaps', 0)}"
                p.level = 0
            
            if tokuro_data.remediation_recommendations:
                p = content.add_paragraph()
                p.text = "Top Recommendations:"
                p.level = 0
                
                for rec in tokuro_data.remediation_recommendations[:3]:
                    p = content.add_paragraph()
                    p.text = f"{rec.get('recommendation', 'N/A')} (Priority: {rec.get('priority', 'N/A')})"
                    p.level = 1
        
        # Save presentation
        prs.save(str(output_path))
        
        return str(output_path)

    async def _generate_docx_report(
        self, 
        config: ReportGeneration, 
        data: Dict[str, Any], 
        template: ReportTemplate
    ) -> str:
        """Generate Word document report"""
        
        filename = f"{config.id}_{config.name.replace(' ', '_')}.docx"
        output_path = self.output_dir / filename
        
        # Create document
        doc = Document()
        
        # Title
        title = doc.add_heading(data['report_info']['name'], 0)
        title.alignment = 1  # Center alignment
        
        # Organization and date
        org_para = doc.add_paragraph()
        org_para.alignment = 1
        org_para.add_run(f"Organization: {data['organization'].get('name', 'N/A')}\n")
        org_para.add_run(f"Generated on: {data['report_info']['generated_at'].strftime('%B %d, %Y')}\n")
        org_para.add_run(f"{self.branding['watermark_text']}").italic = True
        
        # Executive Summary
        doc.add_heading('Executive Summary', level=1)
        if config.tokuro_ai_data and config.tokuro_ai_data.executive_insights:
            for insight in config.tokuro_ai_data.executive_insights:
                doc.add_paragraph(insight, style='List Bullet')
        else:
            doc.add_paragraph('Executive summary will be populated with assessment data.')
        
        # Assessment Overview
        doc.add_heading('Assessment Overview', level=1)
        if data.get('assessment'):
            assessment = data['assessment']
            
            table = doc.add_table(rows=1, cols=2)
            table.style = 'Table Grid'
            
            # Header row
            hdr_cells = table.rows[0].cells
            hdr_cells[0].text = 'Property'
            hdr_cells[1].text = 'Value'
            
            # Add data rows
            properties = [
                ('Title', assessment.get('title', 'N/A')),
                ('Framework', assessment.get('framework', 'N/A')),
                ('Status', assessment.get('status', 'N/A')),
                ('Controls Assessed', str(len(data.get('control_assessments', [])))),
                ('Evidence Files', str(len(data.get('evidence_files', []))))
            ]
            
            for prop, value in properties:
                row_cells = table.add_row().cells
                row_cells[0].text = prop
                row_cells[1].text = value
        
        # Tokuro AI Insights
        if config.tokuro_ai_data:
            doc.add_heading('Tokuro AI Insights', level=1)
            
            tokuro_data = config.tokuro_ai_data
            
            if tokuro_data.gap_analysis:
                gap = tokuro_data.gap_analysis
                doc.add_paragraph(f"Total Gaps Identified: {gap.get('total_gaps', 0)}")
                doc.add_paragraph(f"Critical Gaps: {gap.get('critical_gaps', 0)}")
            
            if tokuro_data.remediation_recommendations:
                doc.add_heading('Recommendations', level=2)
                for i, rec in enumerate(tokuro_data.remediation_recommendations, 1):
                    para = doc.add_paragraph()
                    para.add_run(f"{i}. ").bold = True
                    para.add_run(f"{rec.get('recommendation', 'N/A')} ")
                    para.add_run(f"(Priority: {rec.get('priority', 'N/A')})").italic = True
        
        # Save document
        doc.save(str(output_path))
        
        return str(output_path)

    async def _generate_xlsx_report(
        self, 
        config: ReportGeneration, 
        data: Dict[str, Any], 
        template: ReportTemplate
    ) -> str:
        """Generate Excel spreadsheet report"""
        
        filename = f"{config.id}_{config.name.replace(' ', '_')}.xlsx"
        output_path = self.output_dir / filename
        
        # Create workbook
        wb = Workbook()
        
        # Remove default sheet
        wb.remove(wb.active)
        
        # Summary sheet
        ws_summary = wb.create_sheet("Summary")
        
        # Header styling
        header_font = Font(bold=True, color="FFFFFF")
        header_fill = PatternFill(start_color="1f2937", end_color="1f2937", fill_type="solid")
        
        # Report info
        ws_summary['A1'] = "Report Information"
        ws_summary['A1'].font = header_font
        ws_summary['A1'].fill = header_fill
        
        ws_summary['A2'] = "Report Name:"
        ws_summary['B2'] = data['report_info']['name']
        
        ws_summary['A3'] = "Organization:"
        ws_summary['B3'] = data['organization'].get('name', 'N/A')
        
        ws_summary['A4'] = "Generated On:"
        ws_summary['B4'] = data['report_info']['generated_at'].strftime('%B %d, %Y')
        
        ws_summary['A5'] = "Powered by:"
        ws_summary['B5'] = self.branding['watermark_text']
        
        # Assessment Overview
        if data.get('assessment'):
            ws_summary['A7'] = "Assessment Overview"
            ws_summary['A7'].font = header_font
            ws_summary['A7'].fill = header_fill
            
            assessment = data['assessment']
            row = 8
            
            overview_data = [
                ("Title", assessment.get('title', 'N/A')),
                ("Framework", assessment.get('framework', 'N/A')),
                ("Status", assessment.get('status', 'N/A')),
                ("Controls Assessed", len(data.get('control_assessments', []))),
                ("Evidence Files", len(data.get('evidence_files', [])))
            ]
            
            for label, value in overview_data:
                ws_summary[f'A{row}'] = label
                ws_summary[f'B{row}'] = value
                row += 1
        
        # Control Assessments sheet
        controls = data.get('control_assessments', [])
        if controls:
            ws_controls = wb.create_sheet("Control Assessments")
            
            # Headers
            headers = ['Control ID', 'Status', 'Implementation Score', 'Effectiveness Score', 'Maturity Score', 'Overall Score']
            for col, header in enumerate(headers, 1):
                cell = ws_controls.cell(row=1, column=col)
                cell.value = header
                cell.font = header_font
                cell.fill = header_fill
            
            # Data
            for row, control in enumerate(controls, 2):
                ws_controls.cell(row=row, column=1, value=control.get('control_id', 'N/A'))
                ws_controls.cell(row=row, column=2, value=control.get('status', 'N/A'))
                ws_controls.cell(row=row, column=3, value=control.get('implementation_score', 0))
                ws_controls.cell(row=row, column=4, value=control.get('effectiveness_score', 0))
                ws_controls.cell(row=row, column=5, value=control.get('maturity_score', 0))
                ws_controls.cell(row=row, column=6, value=control.get('overall_score', 0.0))
        
        # Tokuro AI Insights sheet
        if config.tokuro_ai_data:
            ws_ai = wb.create_sheet("Tokuro AI Insights")
            
            tokuro_data = config.tokuro_ai_data
            
            # Gap analysis
            ws_ai['A1'] = "Gap Analysis"
            ws_ai['A1'].font = header_font
            ws_ai['A1'].fill = header_fill
            
            if tokuro_data.gap_analysis:
                gap = tokuro_data.gap_analysis
                ws_ai['A2'] = "Total Gaps:"
                ws_ai['B2'] = gap.get('total_gaps', 0)
                ws_ai['A3'] = "Critical Gaps:"
                ws_ai['B3'] = gap.get('critical_gaps', 0)
            
            # Recommendations
            if tokuro_data.remediation_recommendations:
                ws_ai['A5'] = "Recommendations"
                ws_ai['A5'].font = header_font
                ws_ai['A5'].fill = header_fill
                
                # Headers
                rec_headers = ['Control ID', 'Priority', 'Recommendation', 'Estimated Effort']
                for col, header in enumerate(rec_headers, 1):
                    cell = ws_ai.cell(row=6, column=col)
                    cell.value = header
                    cell.font = header_font
                    cell.fill = header_fill
                
                # Data
                for row, rec in enumerate(tokuro_data.remediation_recommendations, 7):
                    ws_ai.cell(row=row, column=1, value=rec.get('control_id', 'N/A'))
                    ws_ai.cell(row=row, column=2, value=rec.get('priority', 'N/A'))
                    ws_ai.cell(row=row, column=3, value=rec.get('recommendation', 'N/A'))
                    ws_ai.cell(row=row, column=4, value=rec.get('estimated_effort', 'N/A'))
        
        # Auto-adjust column widths
        for ws in wb.worksheets:
            for column in ws.columns:
                max_length = 0
                column_letter = column[0].column_letter
                for cell in column:
                    try:
                        if len(str(cell.value)) > max_length:
                            max_length = len(str(cell.value))
                    except:
                        pass
                adjusted_width = (max_length + 2)
                ws.column_dimensions[column_letter].width = adjusted_width
        
        # Save workbook
        wb.save(str(output_path))
        
        return str(output_path)

    async def _finalize_report(self, config: ReportGeneration, output_path: str):
        """Finalize report generation"""
        
        # Get file info
        file_size = os.path.getsize(output_path)
        filename = os.path.basename(output_path)
        
        # Generate download URL (served by API)
        # Use canonical download route: /api/reports/download/{report_id}
        download_url = f"/api/reports/download/{config.id}"

        # Update report record with completed status and file info
        update_data = {
            'output_filename': filename,
            'output_file_path': str(output_path),
            'file_size_bytes': file_size,
            'download_url': download_url,
            'download_expires_at': (datetime.utcnow() + timedelta(days=30)).isoformat(),
            'generation_completed_at': datetime.utcnow().isoformat(),
            'status': ReportStatus.COMPLETED.value,
            'progress_percentage': 100
        }

        await DatabaseOperations.update_one(
            "report_generations",
            {"id": config.id},
            update_data
        )

    async def _update_report_status(
        self, 
        report_id: str, 
        status: ReportStatus, 
        progress: float = None,
        error_message: str = None
    ):
        """Update report generation status"""
        
        update_data = {
            'status': status.value,
            'updated_at': datetime.utcnow()
        }
        
        if progress is not None:
            update_data['progress_percentage'] = progress
        
        if error_message:
            update_data['error_message'] = error_message
        
        if status == ReportStatus.GENERATING and not progress:
            update_data['generation_started_at'] = datetime.utcnow()
        
        await DatabaseOperations.update_one(
            "report_generations",
            {"id": report_id},
            update_data
        )

    async def _update_report_progress(self, report_id: str, progress: float):
        """Update report generation progress"""
        
        await DatabaseOperations.update_one(
            "report_generations",
            {"id": report_id},
            {
                'progress_percentage': progress,
                'updated_at': datetime.utcnow()
            }
        )

    async def get_report_status(self, report_id: str) -> Optional[Dict[str, Any]]:
        """Get current status of report generation"""
        
        try:
            report = await DatabaseOperations.find_one(
                "report_generations",
                {"id": report_id}
            )
            
            if not report:
                return None
            
            return {
                'id': report['id'],
                'name': report['name'],
                'status': report['status'],
                'progress_percentage': report.get('progress_percentage', 0),
                'error_message': report.get('error_message'),
                'created_at': report['created_at'],
                'generation_started_at': report.get('generation_started_at'),
                'generation_completed_at': report.get('generation_completed_at'),
                'download_url': report.get('download_url'),
                'file_size_bytes': report.get('file_size_bytes')
            }
            
        except Exception as e:
            logger.error(f"Error getting report status: {str(e)}")
            return None

    async def cleanup_old_reports(self, days_old: int = 30):
        """Clean up old report files"""
        
        try:
            cutoff_date = datetime.utcnow() - timedelta(days=days_old)
            
            # Find old reports
            old_reports = await DatabaseOperations.find_many(
                "report_generations",
                {"created_at": {"$lt": cutoff_date}}
            )
            
            cleaned_count = 0
            for report in old_reports:
                # Delete file if exists
                if report.get('output_file_path') and os.path.exists(report['output_file_path']):
                    os.remove(report['output_file_path'])
                
                # Remove from database
                await DatabaseOperations.delete_one(
                    "report_generations",
                    {"id": report['id']}
                )
                
                cleaned_count += 1
            
            logger.info(f"Cleaned up {cleaned_count} old reports")
            return cleaned_count
            
        except Exception as e:
            logger.error(f"Error cleaning up old reports: {str(e)}")
            return 0